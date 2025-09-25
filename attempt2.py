# doc_generator.py
import os
import sys
import subprocess
import json
from openai import OpenAI
from dotenv import load_dotenv

# ==============================================================================
# --- 1. CONFIGURATION & CORE FUNCTIONS ---
# ==============================================================================

# Maximum number of times the script will ask the AI to correct its own errors.
MAX_RETRIES = 3

def setup_client():
    """Loads API key from .env and sets up the OpenAI client."""
    load_dotenv()
    api_key="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", # Replace with your OpenRouter key

    if not api_key:
        print("ERROR: OPENROUTER_API_KEY not found in environment variables.")
        print("Please create a .env file with OPENROUTER_API_KEY='your-key'")
        sys.exit(1)
        
    client = OpenAI(
      base_url="https://openrouter.ai/api/v1",
      api_key="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" # Replace with your OpenRouter key,
    )
    return client

def generate_and_run_script(client, initial_prompt, script_filename, topic):
    """
    Generates, runs, and automatically self-corrects a single Python script.
    """
    current_prompt = initial_prompt
    generated_code = ""
    
    try:
        for attempt in range(MAX_RETRIES):
            print(f"\n--- Attempt {attempt + 1} of {MAX_RETRIES} for '{topic}' ---")
            try:
                print("Requesting code from the AI model...")
                completion = client.chat.completions.create(
                    model="deepseek/deepseek-chat-v3.1:free",
                    messages=[{"role": "user", "content": current_prompt}]
                )
                generated_code = completion.choices[0].message.content.strip()

                if generated_code.startswith("```python"):
                    generated_code = generated_code[9:]
                if generated_code.endswith("```"):
                    generated_code = generated_code[:-3]

                print(f"Code received. Saving to '{script_filename}'...")
                with open(script_filename, "w", encoding="utf-8") as f:
                    f.write(generated_code)

                print(f"Executing '{script_filename}'...")
                result = subprocess.run(
                    [sys.executable, script_filename],
                    capture_output=True, text=True, check=True, encoding="utf-8"
                )
                
                print("\n--- ✅ Script Executed Successfully! ---")
                print("--- Output from Generated Script ---")
                print(result.stdout)
                print("------------------------------------")
                return # Success, exit the function

            except subprocess.CalledProcessError as e:
                print(f"\n--- ❌ Script Execution Failed on Attempt {attempt + 1} ---")
                print("STDOUT:", e.stdout)
                print("STDERR:", e.stderr)
                
                if attempt < MAX_RETRIES - 1:
                    print("Asking AI to self-correct the code...")
                    current_prompt = f"""
The following Python script you generated failed.
Original Goal: {topic}

--- FAILED CODE ---
{generated_code}
--- END FAILED CODE ---

It produced this error:
--- ERROR ---
{e.stderr}
--- END ERROR ---

Please analyze the error and provide a corrected, complete version of the Python script. Your entire output must be ONLY the corrected Python code.
"""
                else:
                    print(f"\n--- 💥 Failed to generate working code after {MAX_RETRIES} attempts. ---")
            except Exception as e:
                print(f"\nAn unexpected error occurred during an attempt: {e}")
                return # Exit on other unexpected errors
    
    finally:
        if os.path.exists(script_filename):
            print(f"Cleaning up by deleting '{script_filename}'...")
            os.remove(script_filename)

def generate_and_build_project(client, initial_prompt, topic):
    """
    Generates, validates, and builds a multi-file project from a JSON blueprint.
    Includes self-correction for invalid JSON.
    """
    current_prompt = initial_prompt
    project_blueprint = None
    generated_text = ""

    for attempt in range(MAX_RETRIES):
        print(f"\n--- Attempt {attempt + 1} of {MAX_RETRIES} for '{topic}' ---")
        try:
            print("Requesting project blueprint from the AI model...")
            completion = client.chat.completions.create(
                model="deepseek/deepseek-chat-v3.1:free",
                messages=[{"role": "user", "content": current_prompt}]
            )
            generated_text = completion.choices[0].message.content.strip()

            print("Blueprint received. Parsing JSON...")
            project_blueprint = json.loads(generated_text)
            print("--- ✅ JSON Blueprint is valid! ---")
            break # Success, exit the loop

        except json.JSONDecodeError as e:
            print(f"\n--- ❌ AI returned invalid JSON on Attempt {attempt + 1} ---")
            print(f"JSON Error: {e}")
            if attempt < MAX_RETRIES - 1:
                print("Asking AI to correct the JSON format...")
                current_prompt = f"""
The JSON blueprint you generated was invalid and could not be parsed.
Original Goal: {topic}

--- INVALID OUTPUT ---
{generated_text}
--- END INVALID OUTPUT ---

The parser failed with this error: {e}

Please provide a corrected version of the JSON blueprint. Ensure it is a single, valid JSON object with no extra text or markdown.
"""
            else:
                 print(f"\n--- 💥 Failed to generate a valid JSON blueprint after {MAX_RETRIES} attempts. ---")
                 return
        except Exception as e:
            print(f"\nAn unexpected error occurred: {e}")
            return

    if not project_blueprint:
        return

    new_project_name = project_blueprint.get("project_name", "generated_project")
    files_to_create = project_blueprint.get("files", [])
    
    if not files_to_create:
        print("Warning: The AI returned a blueprint with no files to create.")
        return

    print(f"Starting to build project: '{new_project_name}'...")
    os.makedirs(new_project_name, exist_ok=True)

    for file_info in files_to_create:
        full_path = os.path.join(new_project_name, file_info.get("path"))
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        print(f"  - Creating file: {full_path}")
        with open(full_path, "w", encoding="utf-8") as f:
            f.write(file_info.get("content", ""))

    print("\n--- ✅ Project Generation Complete! ---")
    print(f"Your project has been created in the '{new_project_name}' directory.")

    if any(f["path"] == 'requirements.txt' for f in files_to_create):
        if input("A 'requirements.txt' was found. Install dependencies? (y/n): ").lower() == 'y':
            req_path = os.path.join(new_project_name, 'requirements.txt')
            subprocess.run([sys.executable, "-m", "pip", "install", "-r", req_path], check=True)
            print("Dependencies installed.")

# ==============================================================================
# --- 2. PROMPT DEFINITION FUNCTIONS (Restored to Full Detail) ---
# ==============================================================================

def get_pdf_generation_prompt():
    """Gets user input and constructs the detailed prompt for PDF generation."""
    print("\n--- PDF Report Generator ---")
    topic = input("Enter the topic for the PDF report (e.g., 'The History of AI'): ").strip()
    image_query = input(f"Enter an image search query for '{topic}' (e.g., 'vintage computer'): ").strip()

    past_prompt = """
import os, requests
from reportlab.platypus import SimpleDocTemplate, Paragraph, Image, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch

def download_image(url, filename):
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        with open(filename, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        print(f"Successfully downloaded {filename}")
        return filename
    except requests.exceptions.RequestException as e:
        print(f"Warning: Could not download image. Error: {e}")
        return None

temp_files = []
try:
    doc = SimpleDocTemplate("example_report.pdf")
    Story = []
    styles = getSampleStyleSheet()
    Story.append(Paragraph("Example Report", styles['h1']))
    image_url = 'https://i.imgur.com/5c0Y18d.png'
    local_image_path = 'temp_image.png'
    downloaded_path = download_image(image_url, local_image_path)
    if downloaded_path:
        temp_files.append(downloaded_path)
        img = Image(downloaded_path, width=4*inch, height=2*inch)
        Story.append(img)
    doc.build(Story)
    print("PDF created successfully.")
finally:
    print("Cleaning up temporary files...")
    for path in temp_files:
        if os.path.exists(path):
            os.remove(path)
            print(f"Deleted {path}")
"""
    final_prompt = f"""
You are an expert Python script generator specializing in creating professional PDF documents using the `reportlab` library.
The user wants a script that creates a detailed, multi-page PDF report on the topic: "{topic}".

You MUST follow ALL of these rules:
1.  Your ENTIRE output must be ONLY the Python code. Do not include any explanations or markdown.
2.  The script must use the `reportlab` and `requests` libraries.
3.  **Content Generation**: The script MUST generate a detailed, multi-section report about "{topic}". The content should be comprehensive, well-structured, and include an introduction, multiple body paragraphs, and a conclusion.
4.  **Image Sourcing & Handling**:
    a. The script MUST find a royalty-free, reliable image URL related to the user's image query: "{image_query}".
    b. It MUST include a `download_image` function to download this image to a temporary local file.
    c. The `Image` object in ReportLab must use this downloaded local file.
5.  **Cleanup**: The script MUST use a `try...finally` block to ensure that the temporarily downloaded image file is always deleted after execution, even if errors occur.
6.  **Advanced Features**: The script should use different Paragraph styles (e.g., 'h1', 'h2', 'BodyText').
7.  **Filename**: Save the final PDF to a file named based on the user's topic (e.g., '{topic.lower().replace(' ', '_')}_report.pdf').

--- EXAMPLE SCRIPT TO MIMIC ---
{past_prompt}
--- END EXAMPLE SCRIPT ---

Now, generate the complete Python script for the user's request: "{topic}".
"""
    return final_prompt, topic

def get_diagram_generation_prompt():
    """Gets user input and constructs the detailed prompt for Graphviz diagram generation."""
    print("\n--- Pipeline Diagram Generator ---")
    topic = input("Enter the topic for the pipeline diagram (e.g., 'Customer Churn Prediction'): ").strip()
    
    past_prompt = """
from graphviz import Digraph
dot = Digraph(comment='ML Pipeline')
dot.attr(rankdir='TB', splines='ortho', nodesep='1.0', ranksep='1.5', fontname='Arial', fontsize='12', label='Machine Learning Pipeline for SEER CRC Dataset', labelloc='t')
dot.attr('node', shape='none', margin='0')
main_node_color = '#4a90e2'
detail_node_color = '#ffffff'
border_color = '#cccccc'
dot.node('preprocessing', label='⚙️ DATA PREPROCESSING', shape='box', style='rounded,filled', fillcolor=main_node_color, fontcolor='white', width='4')
preprocessing_details_html = f'''<
<TABLE BORDER="1" CELLBORDER="0" CELLSPACING="0" CELLPADDING="8" BGCOLOR="{detail_node_color}" COLOR="{border_color}">
  <TR><TD ALIGN="LEFT">
    1. Remove columns with &gt;20% missing values<BR/>
    2. Correlation Grouping (|r|&gt;0.4)<BR/>
    3. Iterative Imputation (RF)<BR/>
  </TD></TR>
</TABLE>>'''
dot.node('preprocessing_details', label=preprocessing_details_html)
dot.edge('dataset', 'preprocessing')
with dot.subgraph() as s:
    s.attr(rank='same')
    s.node('preprocessing')
    s.node('preprocessing_details')
dot.render('example_ml_pipeline', format='png', view=False)
"""
    final_prompt = f"""
You are an expert Python script generator specializing in creating elegant `graphviz` diagrams.
The user wants a script that creates a detailed pipeline diagram for the topic: "{topic}".

You MUST follow ALL of these rules:
1.  Your entire output must be ONLY the Python code. Do NOT include any explanations or markdown formatting.
2.  The script must use the `graphviz` Python library.
3.  The script must generate a diagram with at least 3 main stages relevant to the topic.
4.  The script MUST perfectly replicate the elegant, two-column style shown in the example. This includes a main title, a modern color scheme ('#4a90e2'), and 'Arial' font.
5.  The script MUST use HTML-like labels (<TABLE>...) to structure the text inside the detail nodes, just like the example. Use bold tags (<B>) for headers and bullet points (&#8226;).
6.  The script MUST use `rank='same'` subgraphs to align the main pipeline nodes on the left with their corresponding detail boxes on the right.
7.  CRITICAL: The content of the nodes (the main pipeline stages and the text in the detail boxes) MUST be adapted to be relevant to the user's specific topic: "{topic}".
8.  Save the final diagram to a PNG file named after the user's topic (e.g., '{topic.lower().replace(' ', '_')}_pipeline.png') and set `view=False`.

--- EXAMPLE SCRIPT TO MIMIC ---
{past_prompt}
--- END EXAMPLE SCRIPT ---

Now, generate the complete Python script for the user's request: "{topic}".
"""
    return final_prompt, topic

def get_ppt_generation_prompt():
    """Gets user input and constructs the detailed prompt for PowerPoint generation."""
    print("\n--- PowerPoint Presentation Generator ---")
    topic = input("Enter the topic for the PowerPoint presentation: ").strip()

    past_prompt = """
import os, requests
from pptx import Presentation
from pptx.util import Inches
# This is a conceptual example. The real script will also use Pillow.
temp_files = []
try:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Example Presentation"
    print("Example script: Presentation created.")
finally:
    print("Cleaning up temporary image files...")
    # for path in temp_files:
    #     if os.path.exists(path): os.remove(path)
"""
    final_prompt = f"""
You are an expert Python script generator. Your task is to generate a complete, immediately executable Python script.
The user wants a script that creates a PowerPoint presentation about: "{topic}".

You MUST follow ALL of these rules:
1.  Your entire output must be ONLY the Python code. Do NOT include any explanations or markdown.
2.  The script must use `python-pptx`, `requests`, and `Pillow`.
3.  The presentation must contain at least 4 slides: a title slide and three content slides.
4.  Each content slide should have a title and a mix of bullet points and an image.
5.  The script MUST find its own royalty-free image URLs for the content.
6.  The image download function MUST be resilient. It should use a `try...except` block to catch `requests.exceptions.RequestException`. If an image fails to download, it should print a warning and return `None`. The main script must check if the path is not `None` before adding it to a slide.
7.  The script MUST clean up all temporary image files using a `try...finally` block.
8.  Save the final presentation based on the topic (e.g., '{topic.lower().replace(' ', '_')}.pptx').

--- EXAMPLE SCRIPT ---
{past_prompt}
--- END EXAMPLE SCRIPT ---

Now, generate the complete Python script for the user's request: "{topic}".
"""
    return final_prompt, topic

def get_word_generation_prompt():
    """Gets user input and constructs the detailed prompt for Word document generation."""
    print("\n--- Word Document Generator ---")
    topic = input("Enter the topic for the Word document (e.g., 'Quantum Computing'): ").strip()
    image_query = input(f"Enter an image search query for '{topic}' (e.g., 'quantum computer illustration'): ").strip()
    
    past_prompt = """
import os, requests
from docx import Document
from docx.shared import Inches

def download_image(urls, filename):
    for url in urls:
        try:
            response = requests.get(url, stream=True)
            response.raise_for_status()
            with open(filename, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192): f.write(chunk)
            print(f"Successfully downloaded '{filename}'")
            return filename
        except requests.exceptions.RequestException as e:
            print(f"Warning: Failed to download from {url}: {e}")
    return None

temp_files = []
try:
    image_urls = ['https://i.imgur.com/5c0Y18d.png', 'https://i.imgur.com/another-url.png']
    downloaded_path = download_image(image_urls, 'temp_image.png')
    if downloaded_path: temp_files.append(downloaded_path)
    document = Document()
    document.add_heading('Example Document', level=0)
    document.add_paragraph("This is an example.")
    if downloaded_path: document.add_picture(downloaded_path, width=Inches(4.0))
    document.save('example.docx')
    print("Document saved successfully.")
finally:
    print("--- Starting Cleanup ---")
    for path in temp_files:
        if os.path.exists(path): os.remove(path)
"""
    final_prompt = f"""
You are an expert Python script generator that creates detailed Microsoft Word documents (.docx).
The user wants a Word document about the topic: "{topic}".

You MUST follow ALL of these rules:
1.  Your ENTIRE output must be ONLY the Python code. Do not include any explanations or markdown.
2.  The script must use the `python-docx`, `requests`, and `os` libraries.
3.  **Content Generation**: The script must contain a detailed, multi-paragraph summary of "{topic}". The document must be well-structured with an introduction, at least two main body sections with subheadings (using 'Heading 1' and 'Heading 2' styles), and a conclusion.
4.  **Image Sourcing**: The script MUST find **two** different, royalty-free, and reliable image URLs related to "{image_query}".
5.  **Image Download Logic**: The `download_image` function must try each URL in order and stop after the first successful download.
6.  **Cleanup**: The script MUST use a `try...finally` block to ensure any downloaded image file is deleted at the end.
7.  **Conditional Image Placement**: The script must only add the image to the document if the download was successful.
8.  **Filename**: Save the final Word document to a file named based on the topic (e.g., '{topic.lower().replace(' ', '_')}.docx').
9.  **Compatibility**: Do NOT use any emojis in the generated script's `print` statements.

--- EXAMPLE SCRIPT TO MIMIC ---
{past_prompt}
--- END EXAMPLE SCRIPT ---

Now, generate the complete Python script for the user's request.
"""
    return final_prompt, topic

def get_project_generation_prompt():
    """Gets user input and constructs the prompt for a complex code project."""
    print("\n--- Complex Code Project Generator ---")
    topic = input("Describe the project you want to build (e.g., 'A simple FastAPI server with one endpoint'): ").strip()

    past_prompt = r"""
{
  "project_name": "basic_flask_app",
  "files": [
    {
      "path": "app.py",
      "content": "from flask import Flask\n\napp = Flask(__name__)\n\n@app.route('/')\ndef hello_world():\n    return 'Hello, World!'\n\nif __name__ == '__main__':\n    app.run(debug=True)\n"
    },
    {
      "path": "requirements.txt",
      "content": "Flask\n"
    }
  ]
}
"""
    final_prompt = f"""
You are an expert software architect. Your task is to generate a complete, multi-file project as a single JSON object.

RULES:
1. Your ENTIRE output must be ONLY a single JSON object. No explanations or markdown.
2. The JSON must have keys "project_name" (string) and "files" (a list of file objects).
3. Each file object in the "files" list must have keys "path" (string, e.g., "src/main.py") and "content" (string, the full code for the file).
4. If dependencies are needed, you MUST include a "requirements.txt" file.
5. Ensure file content is properly escaped for JSON (e.g., newlines as \\n).

--- EXAMPLE BLUEPRINT TO MIMIC ---
{past_prompt}
--- END EXAMPLE BLUEPRINT ---

Now, generate the complete JSON blueprint for this project: "{topic}"
"""
    return final_prompt, topic

# ==============================================================================
# --- 3. MAIN EXECUTION LOGIC ---
# ==============================================================================

def main():
    """Main function to display the menu and run the document generator."""
    client = setup_client()

    while True:
        print("\n=============================================")
        print("    AI-Powered Document & Code Generator   ")
        print("=============================================")
        print("Select an option:")
        print("  1. PDF Report")
        print("  2. Pipeline Diagram (PNG)")
        print("  3. PowerPoint Presentation (PPTX)")
        print("  4. Word Document (DOCX)")
        print("  5. Complex Code Project")
        print("  6. Exit")
        
        choice = input("\nEnter your choice (1-6): ").strip()
        
        final_prompt = None
        topic = None
        
        if choice in ['1', '2', '3', '4']:
            script_filename = "generated_temp_script.py"
            if choice == '1':
                final_prompt, topic = get_pdf_generation_prompt()
            elif choice == '2':
                final_prompt, topic = get_diagram_generation_prompt()
            elif choice == '3':
                final_prompt, topic = get_ppt_generation_prompt()
            elif choice == '4':
                final_prompt, topic = get_word_generation_prompt()
            
            if final_prompt and topic:
                generate_and_run_script(client, final_prompt, script_filename, topic)

        elif choice == '5':
            final_prompt, topic = get_project_generation_prompt()
            if final_prompt and topic:
                generate_and_build_project(client, final_prompt, topic)

        elif choice == '6':
            print("Exiting program.")
            break
        else:
            print("Invalid choice. Please enter a number between 1 and 6.")

if __name__ == "__main__":
    main()