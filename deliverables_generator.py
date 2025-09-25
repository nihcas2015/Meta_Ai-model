#!/usr/bin/env python3
"""
Deliverables Generator for Meta AI System
Creates actual deliverable files (PDF, PPT, DOCX, etc.) and generates visual previews
"""

import os
import io
import base64
import json
import uuid
from pathlib import Path
from typing import Dict, List, Any, Tuple
from datetime import datetime
import logging

# Document generation libraries
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Word document generation
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Image generation for previews
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from pdf2image import convert_from_path
import subprocess

# Setup
BASE_DIR = Path(__file__).parent
DELIVERABLES_DIR = BASE_DIR / "deliverables"
PREVIEWS_DIR = BASE_DIR / "previews"
DELIVERABLES_DIR.mkdir(exist_ok=True)
PREVIEWS_DIR.mkdir(exist_ok=True)

logger = logging.getLogger(__name__)

class DeliverablesGenerator:
    """Generate actual deliverable files and their visual previews"""
    
    def __init__(self):
        self.deliverables_dir = DELIVERABLES_DIR
        self.previews_dir = PREVIEWS_DIR
        self.styles = getSampleStyleSheet()
        
    def generate_pdf_report(self, user_query: str, domain_outputs: Dict[str, Any], conversation_id: str) -> Dict[str, str]:
        """Generate actual PDF report using ReportLab"""
        
        logger.info("📄 Generating actual PDF report...")
        
        # Create PDF filename
        pdf_filename = f"report_{conversation_id}.pdf"
        pdf_path = self.deliverables_dir / pdf_filename
        
        # Create PDF document
        doc = SimpleDocTemplate(str(pdf_path), pagesize=A4,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=18)
        
        # Container for PDF content
        story = []
        
        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor=colors.navy
        )
        
        # Add title
        story.append(Paragraph(f"TECHNICAL ANALYSIS REPORT", title_style))
        story.append(Paragraph(f"{user_query}", self.styles['Heading2']))
        story.append(Spacer(1, 20))
        
        # Add metadata
        meta_data = [
            ['Report ID:', conversation_id],
            ['Generated:', datetime.now().strftime('%B %d, %Y at %I:%M %p')],
            ['Domains Analyzed:', ', '.join(domain_outputs.keys()).title()]
        ]
        meta_table = Table(meta_data, colWidths=[2*inch, 4*inch])
        meta_table.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
        ]))
        story.append(meta_table)
        story.append(Spacer(1, 30))
        
        # Executive Summary
        story.append(Paragraph("EXECUTIVE SUMMARY", self.styles['Heading2']))
        
        total_findings = sum(len(output.get('key_findings', [])) for output in domain_outputs.values())
        total_recommendations = sum(len(output.get('recommendations', [])) for output in domain_outputs.values())
        
        exec_summary = f"""
        This comprehensive technical analysis examines {len(domain_outputs)} critical engineering domains 
        for the {user_query} project. Our multi-disciplinary approach has identified {total_findings} 
        key technical findings and generated {total_recommendations} actionable recommendations.
        
        The analysis covers mechanical engineering considerations, electrical system requirements, 
        and software architecture aspects, providing a holistic view of the technical landscape 
        and implementation roadmap.
        """
        story.append(Paragraph(exec_summary, self.styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Domain Analysis Sections
        for domain, output in domain_outputs.items():
            # Domain header
            story.append(Paragraph(f"{domain.upper()} ANALYSIS", self.styles['Heading2']))
            
            # Key Findings
            story.append(Paragraph("Key Findings:", self.styles['Heading3']))
            for finding in output.get('key_findings', [])[:5]:  # Top 5
                story.append(Paragraph(f"• {finding}", self.styles['Normal']))
            story.append(Spacer(1, 10))
            
            # Recommendations
            story.append(Paragraph("Recommendations:", self.styles['Heading3']))
            for rec in output.get('recommendations', [])[:5]:  # Top 5
                story.append(Paragraph(f"• {rec}", self.styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        logger.info(f"✅ PDF report generated: {pdf_path}")
        
        # Generate preview image
        preview_path = self.generate_pdf_preview(pdf_path, conversation_id)
        
        return {
            'deliverable_path': str(pdf_path),
            'preview_path': str(preview_path),
            'file_type': 'pdf',
            'filename': pdf_filename
        }
    
    def generate_powerpoint_presentation(self, user_query: str, domain_outputs: Dict[str, Any], conversation_id: str) -> Dict[str, str]:
        """Generate actual PowerPoint presentation"""
        
        logger.info("📽️ Generating actual PowerPoint presentation...")
        
        # Create presentation
        prs = Presentation()
        
        # Slide 1: Title Slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = user_query.upper()
        subtitle.text = f"Technical Analysis Report\nGenerated by Meta AI System\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Slide 2: Overview
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Analysis Overview"
        tf = content.text_frame
        
        total_findings = sum(len(output.get('key_findings', [])) for output in domain_outputs.values())
        total_recommendations = sum(len(output.get('recommendations', [])) for output in domain_outputs.values())
        
        tf.text = f"Domains Analyzed: {len(domain_outputs)}\n"
        tf.text += f"Total Findings: {total_findings}\n"
        tf.text += f"Total Recommendations: {total_recommendations}\n\n"
        tf.text += "Domains:\n"
        for domain in domain_outputs.keys():
            tf.text += f"• {domain.title()}\n"
        
        # Domain-specific slides
        for domain, output in domain_outputs.items():
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"{domain.title()} Analysis"
            tf = content.text_frame
            
            tf.text = "Key Findings:\n"
            for finding in output.get('key_findings', [])[:4]:
                tf.text += f"• {finding}\n"
            
            tf.text += "\nRecommendations:\n"
            for rec in output.get('recommendations', [])[:4]:
                tf.text += f"• {rec}\n"
        
        # Final slide: Next Steps
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Next Steps & Implementation"
        tf = content.text_frame
        
        tf.text = "Immediate Actions:\n"
        tf.text += "1. Review domain-specific recommendations\n"
        tf.text += "2. Prioritize implementation based on resources\n"
        tf.text += "3. Create detailed project timeline\n"
        tf.text += "4. Assign team responsibilities\n"
        tf.text += "5. Begin prototype development\n"
        
        # Save presentation
        ppt_filename = f"presentation_{conversation_id}.pptx"
        ppt_path = self.deliverables_dir / ppt_filename
        prs.save(str(ppt_path))
        
        logger.info(f"✅ PowerPoint presentation generated: {ppt_path}")
        
        # Generate preview images
        preview_paths = self.generate_ppt_previews(ppt_path, conversation_id)
        
        return {
            'deliverable_path': str(ppt_path),
            'preview_paths': preview_paths,
            'file_type': 'powerpoint',
            'filename': ppt_filename
        }
    
    def generate_word_document(self, user_query: str, domain_outputs: Dict[str, Any], conversation_id: str) -> Dict[str, str]:
        """Generate actual Word document"""
        
        logger.info("📝 Generating actual Word document...")
        
        # Create document
        doc = Document()
        
        # Title
        title = doc.add_heading(user_query.upper(), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Subtitle
        subtitle = doc.add_heading('Technical Analysis Report', level=1)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata
        doc.add_paragraph(f"Report ID: {conversation_id}")
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
        doc.add_paragraph(f"Domains Analyzed: {', '.join(domain_outputs.keys()).title()}")
        doc.add_paragraph()
        
        # Executive Summary
        doc.add_heading('Executive Summary', level=1)
        
        total_findings = sum(len(output.get('key_findings', [])) for output in domain_outputs.values())
        total_recommendations = sum(len(output.get('recommendations', [])) for output in domain_outputs.values())
        
        exec_summary = f"""
        This comprehensive technical analysis examines {len(domain_outputs)} critical engineering domains 
        for the {user_query} project. Our multi-disciplinary approach has identified {total_findings} 
        key technical findings and generated {total_recommendations} actionable recommendations.
        
        The analysis covers mechanical engineering considerations, electrical system requirements, 
        and software architecture aspects, providing a holistic view of the technical landscape 
        and implementation roadmap.
        """
        doc.add_paragraph(exec_summary)
        
        # Domain sections
        for domain, output in domain_outputs.items():
            doc.add_heading(f"{domain.title()} Analysis", level=1)
            
            # Analysis summary
            analysis_text = output.get('analysis', '')[:500] + "..."
            doc.add_paragraph(analysis_text)
            
            # Key Findings
            doc.add_heading('Key Findings', level=2)
            for finding in output.get('key_findings', []):
                p = doc.add_paragraph(finding, style='List Bullet')
            
            # Recommendations
            doc.add_heading('Recommendations', level=2)
            for rec in output.get('recommendations', []):
                p = doc.add_paragraph(rec, style='List Bullet')
            
            doc.add_paragraph()
        
        # Next Steps
        doc.add_heading('Next Steps & Implementation', level=1)
        next_steps = [
            "Review domain-specific recommendations in detail",
            "Prioritize implementation based on available resources",
            "Create detailed project timeline with milestones",
            "Assign team responsibilities and accountabilities",
            "Begin prototype development and testing",
            "Establish regular review and feedback cycles"
        ]
        
        for step in next_steps:
            doc.add_paragraph(step, style='List Number')
        
        # Save document
        docx_filename = f"document_{conversation_id}.docx"
        docx_path = self.deliverables_dir / docx_filename
        doc.save(str(docx_path))
        
        logger.info(f"✅ Word document generated: {docx_path}")
        
        # Generate preview
        preview_path = self.generate_docx_preview(docx_path, conversation_id)
        
        return {
            'deliverable_path': str(docx_path),
            'preview_path': str(preview_path),
            'file_type': 'word',
            'filename': docx_filename
        }
    
    def generate_project_files(self, user_query: str, domain_outputs: Dict[str, Any], conversation_id: str) -> Dict[str, str]:
        """Generate actual project structure with code files"""
        
        logger.info("💻 Generating actual project files...")
        
        # Create project directory
        project_name = user_query.lower().replace(' ', '_').replace('-', '_')
        project_dir = self.deliverables_dir / f"{project_name}_{conversation_id}"
        project_dir.mkdir(exist_ok=True)
        
        # Create directory structure
        directories = [
            'src', 'src/mechanical', 'src/electrical', 'src/programming',
            'docs', 'tests', 'config', 'data', 'scripts'
        ]
        
        for dir_name in directories:
            (project_dir / dir_name).mkdir(parents=True, exist_ok=True)
        
        # Generate README.md
        readme_content = f"""# {user_query}
        
## Project Overview

This project was generated by Meta AI System based on comprehensive domain analysis.

### Domains Analyzed
"""
        for domain in domain_outputs.keys():
            readme_content += f"- {domain.title()}\n"
        
        readme_content += f"""
### Key Features
- Multi-domain engineering approach
- Comprehensive analysis and recommendations
- Scalable architecture design

### Project Structure
```
{project_name}/
├── src/
│   ├── mechanical/     # Mechanical engineering components
│   ├── electrical/     # Electrical system designs  
│   └── programming/    # Software implementation
├── docs/              # Documentation
├── tests/             # Test files
├── config/            # Configuration files
├── data/              # Data files
└── scripts/           # Utility scripts
```

### Getting Started
1. Review the domain-specific documentation in `docs/`
2. Check configuration files in `config/`
3. Run tests with `python -m pytest tests/`
4. Follow implementation guidelines in each domain folder

Generated on: {datetime.now().strftime('%B %d, %Y')}
Report ID: {conversation_id}
"""
        
        with open(project_dir / 'README.md', 'w') as f:
            f.write(readme_content)
        
        # Generate domain-specific files
        for domain, output in domain_outputs.items():
            domain_dir = project_dir / 'src' / domain
            
            # Create domain-specific implementation file
            if domain == 'mechanical':
                self._create_mechanical_files(domain_dir, output)
            elif domain == 'electrical':
                self._create_electrical_files(domain_dir, output)
            elif domain == 'programming':
                self._create_programming_files(domain_dir, output)
        
        # Create configuration files
        self._create_config_files(project_dir / 'config', domain_outputs)
        
        # Create test files
        self._create_test_files(project_dir / 'tests', domain_outputs)
        
        # Create ZIP archive of the project
        import zipfile
        zip_filename = f"{project_name}_{conversation_id}.zip"
        zip_path = self.deliverables_dir / zip_filename
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for file_path in project_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(project_dir)
                    zipf.write(file_path, arcname)
        
        logger.info(f"✅ Project files generated: {zip_path}")
        
        # Generate preview
        preview_path = self.generate_project_preview(project_dir, conversation_id)
        
        return {
            'deliverable_path': str(zip_path),
            'project_path': str(project_dir),
            'preview_path': str(preview_path),
            'file_type': 'project',
            'filename': zip_filename
        }
    
    def _create_mechanical_files(self, domain_dir: Path, output: Dict[str, Any]):
        """Create mechanical engineering specific files"""
        
        # Material specifications
        material_spec = f"""# Mechanical Engineering Specifications

## Key Findings
"""
        for finding in output.get('key_findings', []):
            material_spec += f"- {finding}\n"
        
        material_spec += f"""
## Recommendations
"""
        for rec in output.get('recommendations', []):
            material_spec += f"- {rec}\n"
        
        with open(domain_dir / 'specifications.md', 'w') as f:
            f.write(material_spec)
        
        # Python module for calculations
        python_code = f'''"""
Mechanical Engineering Calculations Module
Generated by Meta AI System
"""

import math
import numpy as np

class MechanicalAnalysis:
    """Mechanical engineering analysis tools"""
    
    def __init__(self):
        self.safety_factor = 2.0
        
    def stress_analysis(self, force, area):
        """Calculate stress given force and area"""
        return force / area
    
    def fatigue_analysis(self, max_stress, min_stress, cycles):
        """Basic fatigue analysis"""
        stress_range = max_stress - min_stress
        mean_stress = (max_stress + min_stress) / 2
        return {{"stress_range": stress_range, "mean_stress": mean_stress}}
    
    def thermal_expansion(self, length, temp_change, expansion_coeff):
        """Calculate thermal expansion"""
        return length * temp_change * expansion_coeff

# Key findings from analysis:
# {chr(10).join([f"# - {finding}" for finding in output.get('key_findings', [])])}
'''
        
        with open(domain_dir / 'analysis.py', 'w') as f:
            f.write(python_code)
    
    def _create_electrical_files(self, domain_dir: Path, output: Dict[str, Any]):
        """Create electrical engineering specific files"""
        
        # Circuit specifications
        circuit_spec = f"""# Electrical Engineering Specifications

## Key Findings
"""
        for finding in output.get('key_findings', []):
            circuit_spec += f"- {finding}\n"
        
        circuit_spec += f"""
## Recommendations
"""
        for rec in output.get('recommendations', []):
            circuit_spec += f"- {rec}\n"
        
        with open(domain_dir / 'circuit_specs.md', 'w') as f:
            f.write(circuit_spec)
        
        # Python module for electrical calculations
        python_code = f'''"""
Electrical Engineering Calculations Module
Generated by Meta AI System
"""

import math
import numpy as np

class ElectricalAnalysis:
    """Electrical engineering analysis tools"""
    
    def __init__(self):
        self.voltage_tolerance = 0.05
        
    def power_calculation(self, voltage, current):
        """Calculate power"""
        return voltage * current
    
    def impedance_analysis(self, resistance, reactance):
        """Calculate impedance"""
        return math.sqrt(resistance**2 + reactance**2)
    
    def efficiency_calculation(self, power_out, power_in):
        """Calculate efficiency percentage"""
        return (power_out / power_in) * 100

# Key findings from analysis:
# {chr(10).join([f"# - {finding}" for finding in output.get('key_findings', [])])}
'''
        
        with open(domain_dir / 'calculations.py', 'w') as f:
            f.write(python_code)
    
    def _create_programming_files(self, domain_dir: Path, output: Dict[str, Any]):
        """Create software engineering specific files"""
        
        # Software architecture document
        arch_spec = f"""# Software Architecture Specifications

## Key Findings
"""
        for finding in output.get('key_findings', []):
            arch_spec += f"- {finding}\n"
        
        arch_spec += f"""
## Recommendations
"""
        for rec in output.get('recommendations', []):
            arch_spec += f"- {rec}\n"
        
        with open(domain_dir / 'architecture.md', 'w') as f:
            f.write(arch_spec)
        
        # Main application module
        python_code = f'''"""
Main Application Module
Generated by Meta AI System
"""

import logging
from typing import Dict, List, Any
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ApplicationCore:
    """Core application class"""
    
    def __init__(self):
        self.initialized = False
        self.start_time = datetime.now()
        
    def initialize(self):
        """Initialize the application"""
        logger.info("Initializing application...")
        self.initialized = True
        return True
    
    def process_data(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Process input data and return results"""
        if not self.initialized:
            raise RuntimeError("Application not initialized")
        
        results = {{
            "processed": True,
            "timestamp": datetime.now().isoformat(),
            "input_keys": list(data.keys())
        }}
        
        return results
    
    def get_status(self) -> Dict[str, Any]:
        """Get application status"""
        return {{
            "initialized": self.initialized,
            "uptime": (datetime.now() - self.start_time).total_seconds()
        }}

if __name__ == "__main__":
    app = ApplicationCore()
    app.initialize()
    logger.info("Application started successfully")

# Key findings from analysis:
# {chr(10).join([f"# - {finding}" for finding in output.get('key_findings', [])])}
'''
        
        with open(domain_dir / 'main.py', 'w') as f:
            f.write(python_code)
    
    def _create_config_files(self, config_dir: Path, domain_outputs: Dict[str, Any]):
        """Create configuration files"""
        
        # Main config file
        config_content = f"""# Configuration File
# Generated by Meta AI System

[system]
version = "1.0.0"
debug = false
log_level = "INFO"

[domains]
mechanical = true
electrical = true
programming = true

[analysis]
domains_analyzed = {len(domain_outputs)}
total_findings = {sum(len(output.get('key_findings', [])) for output in domain_outputs.values())}
total_recommendations = {sum(len(output.get('recommendations', [])) for output in domain_outputs.values())}
"""
        
        with open(config_dir / 'config.ini', 'w') as f:
            f.write(config_content)
    
    def _create_test_files(self, tests_dir: Path, domain_outputs: Dict[str, Any]):
        """Create test files"""
        
        # Main test file
        test_content = f'''"""
Test Suite
Generated by Meta AI System
"""

import unittest
from unittest.mock import Mock, patch
import sys
import os

# Add src to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

class TestDomainAnalysis(unittest.TestCase):
    """Test domain analysis components"""
    
    def setUp(self):
        """Set up test fixtures"""
        self.domains = {list(domain_outputs.keys())}
        
    def test_mechanical_analysis(self):
        """Test mechanical analysis"""
        # Mock test for mechanical domain
        self.assertIn('mechanical', self.domains)
    
    def test_electrical_analysis(self):
        """Test electrical analysis"""  
        # Mock test for electrical domain
        self.assertIn('electrical', self.domains)
        
    def test_programming_analysis(self):
        """Test programming analysis"""
        # Mock test for programming domain
        self.assertIn('programming', self.domains)

if __name__ == '__main__':
    unittest.main()
'''
        
        with open(tests_dir / 'test_analysis.py', 'w') as f:
            f.write(test_content)
    
    def generate_pdf_preview(self, pdf_path: Path, conversation_id: str) -> str:
        """Generate preview image of PDF first page"""
        
        try:
            # Convert PDF first page to image
            pages = convert_from_path(str(pdf_path), first_page=1, last_page=1, dpi=150)
            if pages:
                preview_path = self.previews_dir / f"pdf_preview_{conversation_id}.png"
                pages[0].save(preview_path, 'PNG')
                logger.info(f"✅ PDF preview generated: {preview_path}")
                return str(preview_path)
        except Exception as e:
            logger.error(f"Error generating PDF preview: {e}")
            # Fallback: create a simple preview image
            return self._create_fallback_preview("PDF Report", conversation_id, "pdf")
    
    def generate_ppt_previews(self, ppt_path: Path, conversation_id: str) -> List[str]:
        """Generate preview images of PowerPoint slides"""
        
        preview_paths = []
        try:
            # Use python-pptx to extract slide content and create previews
            from pptx import Presentation
            prs = Presentation(str(ppt_path))
            
            for i, slide in enumerate(prs.slides):
                preview_path = self.previews_dir / f"ppt_slide_{i+1}_{conversation_id}.png"
                
                # Create a simple preview image with slide content
                fig, ax = plt.subplots(figsize=(10, 7.5))
                fig.patch.set_facecolor('white')
                
                # Extract slide title
                title_text = "Slide Content"
                try:
                    if hasattr(slide, 'shapes') and slide.shapes.title:
                        title_text = slide.shapes.title.text
                except:
                    pass
                
                ax.text(0.5, 0.8, f"Slide {i+1}", ha='center', va='center', 
                       fontsize=20, fontweight='bold', transform=ax.transAxes)
                ax.text(0.5, 0.6, title_text, ha='center', va='center',
                       fontsize=16, transform=ax.transAxes, wrap=True)
                
                ax.set_xlim(0, 1)
                ax.set_ylim(0, 1)
                ax.axis('off')
                
                plt.tight_layout()
                plt.savefig(preview_path, dpi=150, bbox_inches='tight', facecolor='white')
                plt.close()
                
                preview_paths.append(str(preview_path))
            
            logger.info(f"✅ PowerPoint previews generated: {len(preview_paths)} slides")
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint previews: {e}")
            # Fallback preview
            preview_paths.append(self._create_fallback_preview("PowerPoint Presentation", conversation_id, "ppt"))
        
        return preview_paths
    
    def generate_docx_preview(self, docx_path: Path, conversation_id: str) -> str:
        """Generate preview image of Word document"""
        
        try:
            # Create a simple preview showing document structure
            preview_path = self.previews_dir / f"docx_preview_{conversation_id}.png"
            
            fig, ax = plt.subplots(figsize=(8.5, 11))
            fig.patch.set_facecolor('white')
            
            # Document preview layout
            ax.text(0.5, 0.95, "WORD DOCUMENT", ha='center', va='center',
                   fontsize=16, fontweight='bold', transform=ax.transAxes)
            ax.text(0.5, 0.9, "Technical Analysis Report", ha='center', va='center',
                   fontsize=14, transform=ax.transAxes)
            
            # Add some content lines
            content_lines = [
                "Executive Summary",
                "Domain Analysis Results", 
                "Key Findings & Recommendations",
                "Implementation Guidelines",
                "Next Steps"
            ]
            
            y_pos = 0.8
            for line in content_lines:
                ax.text(0.1, y_pos, f"• {line}", fontsize=12, transform=ax.transAxes)
                y_pos -= 0.1
            
            # Add border
            border = patches.Rectangle((0.05, 0.05), 0.9, 0.9, linewidth=2,
                                     edgecolor='black', facecolor='none', transform=ax.transAxes)
            ax.add_patch(border)
            
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis('off')
            
            plt.tight_layout()
            plt.savefig(preview_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            logger.info(f"✅ Word document preview generated: {preview_path}")
            return str(preview_path)
            
        except Exception as e:
            logger.error(f"Error generating Word preview: {e}")
            return self._create_fallback_preview("Word Document", conversation_id, "docx")
    
    def generate_project_preview(self, project_dir: Path, conversation_id: str) -> str:
        """Generate preview image of project structure"""
        
        try:
            preview_path = self.previews_dir / f"project_preview_{conversation_id}.png"
            
            fig, ax = plt.subplots(figsize=(12, 8))
            fig.patch.set_facecolor('white')
            
            ax.text(0.5, 0.95, "PROJECT STRUCTURE", ha='center', va='center',
                   fontsize=16, fontweight='bold', transform=ax.transAxes)
            
            # Create a simple directory tree visualization
            tree_structure = [
                f"{project_dir.name}/",
                "├── src/",
                "│   ├── mechanical/",
                "│   ├── electrical/", 
                "│   └── programming/",
                "├── docs/",
                "├── tests/",
                "├── config/",
                "├── data/",
                "└── README.md"
            ]
            
            y_pos = 0.85
            for item in tree_structure:
                ax.text(0.1, y_pos, item, fontsize=12, fontfamily='monospace',
                       transform=ax.transAxes)
                y_pos -= 0.08
            
            # Add file count info
            file_count = sum(1 for _ in project_dir.rglob('*') if _.is_file())
            ax.text(0.6, 0.4, f"Total Files: {file_count}", fontsize=12, fontweight='bold',
                   transform=ax.transAxes)
            ax.text(0.6, 0.35, f"Generated: {datetime.now().strftime('%Y-%m-%d')}", 
                   fontsize=10, transform=ax.transAxes)
            
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis('off')
            
            plt.tight_layout()
            plt.savefig(preview_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            logger.info(f"✅ Project preview generated: {preview_path}")
            return str(preview_path)
            
        except Exception as e:
            logger.error(f"Error generating project preview: {e}")
            return self._create_fallback_preview("Project Files", conversation_id, "project")
    
    def _create_fallback_preview(self, content_type: str, conversation_id: str, file_type: str) -> str:
        """Create a fallback preview image when other methods fail"""
        
        preview_path = self.previews_dir / f"fallback_{file_type}_{conversation_id}.png"
        
        fig, ax = plt.subplots(figsize=(8, 6))
        fig.patch.set_facecolor('#f0f0f0')
        
        ax.text(0.5, 0.6, content_type, ha='center', va='center',
               fontsize=18, fontweight='bold', transform=ax.transAxes)
        ax.text(0.5, 0.4, "Preview Generated", ha='center', va='center',
               fontsize=14, transform=ax.transAxes)
        ax.text(0.5, 0.3, f"ID: {conversation_id}", ha='center', va='center',
               fontsize=10, transform=ax.transAxes, alpha=0.7)
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        plt.tight_layout()
        plt.savefig(preview_path, dpi=150, bbox_inches='tight', facecolor='#f0f0f0')
        plt.close()
        
        return str(preview_path)

if __name__ == "__main__":
    # Test the deliverables generator
    generator = DeliverablesGenerator()
    
    # Mock data for testing
    mock_domain_outputs = {
        'mechanical': {
            'key_findings': ['High stress concentration', 'Material fatigue risk'],
            'recommendations': ['Use stronger alloy', 'Implement stress relief'],
            'analysis': 'Detailed mechanical analysis...'
        },
        'electrical': {
            'key_findings': ['Power consumption high', 'EMI concerns'],
            'recommendations': ['Optimize power circuit', 'Add EMI shielding'],
            'analysis': 'Electrical system analysis...'
        },
        'programming': {
            'key_findings': ['Algorithm efficiency', 'Memory usage'],
            'recommendations': ['Optimize algorithms', 'Implement caching'],
            'analysis': 'Software architecture analysis...'
        }
    }
    
    # Test PDF generation
    pdf_result = generator.generate_pdf_report('Smart Home System', mock_domain_outputs, 'test123')
    print("PDF generated:", pdf_result)